#!/usr/bin/env python3
"""
每日PDF处理核心类 - 完整的每日PDF处理流程
包含：PDF识别、内容提取、AI分析和报告生成
"""

import os
import sys
import re
import time
import json
import logging
from datetime import datetime
import subprocess
import shutil
import tempfile

from pdf_ocr_tool.pipeline import pdf_processing_pipeline
from pdf_ocr_tool.summarizers import financial_summarizer as ai_content_summarizer


def copy_reports_to_collections(date_dir):
    """把日期目录 reports/ 下的汇总文档复制到 daily 根目录的集合文件夹

    参考 topics/extractor.py 中"核心论点汇总"的复制规则：
    将文件用 shutil.copy2 复制到固定根目录下的文件夹，扁平结构 + 保留日期后缀命名。
      - 每日重点汇总_*.md -> 重点汇总/
      - summary_list_*.md -> 一句话总结/

    Args:
        date_dir: 日期目录，如 output/daily/20260823

    Returns:
        list[(src, dst)]: 成功复制的 (源文件, 目标文件) 列表
    """
    reports_dir = os.path.join(date_dir, 'reports')
    if not os.path.isdir(reports_dir):
        return []

    daily_root = os.path.dirname(os.path.abspath(date_dir))
    copies = []

    for fname in sorted(os.listdir(reports_dir)):
        if not fname.endswith('.md'):
            continue
        src = os.path.join(reports_dir, fname)
        if fname.startswith('每日重点汇总_'):
            dst_dir = os.path.join(daily_root, '重点汇总')
        elif fname.startswith('summary_list_'):
            dst_dir = os.path.join(daily_root, '一句话总结')
        else:
            continue
        os.makedirs(dst_dir, exist_ok=True)
        dst = os.path.join(dst_dir, fname)
        shutil.copy2(src, dst)
        copies.append((src, dst))

    return copies


def collect_daily_summaries(date_dir):
    """把当日汇总文档收集到 output/汇总/YYYYMMDD/（微信读书导入用）

    将分散在多个目录的当日汇总文档统一复制到 output/汇总/YYYYMMDD/，
    按日期区分，方便整体导入微信读书。收集内容：
      - 每日重点汇总_YYYYMMDD.md    （当日重点总结）
      - summary_list_YYYYMMDD.md    （一句话总结）
      - {专题}_核心论点汇总_YYYYMMDD.md（所有已提取专题的核心论点汇总，
        遍历 output/topic_summaries/*/核心论点/ 自动收集）

    Args:
        date_dir: 日期目录，如 output/daily/20260823

    Returns:
        list[(src, dst)]: 成功复制的 (源文件, 目标文件) 列表
    """
    daily_root = os.path.dirname(os.path.abspath(date_dir))
    output_root = os.path.dirname(os.path.abspath(daily_root))
    date_str = os.path.basename(os.path.normpath(date_dir))

    target_dir = os.path.join(output_root, '汇总', date_str)
    os.makedirs(target_dir, exist_ok=True)

    candidates = []

    reports_dir = os.path.join(date_dir, 'reports')
    if os.path.isdir(reports_dir):
        for fname in os.listdir(reports_dir):
            if not fname.endswith('.md'):
                continue
            if fname.startswith('每日重点汇总_') or fname.startswith('summary_list_'):
                candidates.append(os.path.join(reports_dir, fname))

    topics_root = os.path.join(output_root, 'topic_summaries')
    if os.path.isdir(topics_root):
        for topic_name in sorted(os.listdir(topics_root)):
            core_dir = os.path.join(topics_root, topic_name, '核心论点')
            if not os.path.isdir(core_dir):
                continue
            for fname in os.listdir(core_dir):
                if fname.endswith('.md') and fname.endswith(f'核心论点汇总_{date_str}.md'):
                    candidates.append(os.path.join(core_dir, fname))

    copies = []
    seen = set()
    for src in sorted(candidates):
        src_abs = os.path.abspath(src)
        if src_abs in seen:
            continue
        seen.add(src_abs)
        dst = os.path.join(target_dir, os.path.basename(src))
        shutil.copy2(src, dst)
        copies.append((src, dst))

    return copies


class DailyPDFProcessor:
    """每日PDF处理系统类"""
    
    def __init__(self, config):
        self.config = config
        self.setup_logging()
        
        # 创建临时目录
        self.temp_dir = tempfile.mkdtemp(prefix="pdf_processing_")
        self.logger.info(f"临时目录: {self.temp_dir}")
        
        # 固定缓存位置（在输出根目录）
        self.cache_file = self._get_cache_file()
    
    def _get_cache_file(self):
        """获取固定的缓存文件路径"""
        # 找到 output 根目录
        output_dir = self.config['output_dir']
        # 向上查找直到找到 daily 目录的父目录（output）
        while True:
            parent = os.path.dirname(output_dir)
            if os.path.basename(output_dir) == 'daily' or parent == output_dir:
                break
            output_dir = parent
        
        # 缓存文件放在 output 根目录
        return os.path.join(output_dir, '.pdf_parse_cache.json')
    
    def setup_logging(self):
        """设置日志系统"""
        log_dir = os.path.join(self.config['output_dir'], 'logs')
        os.makedirs(log_dir, exist_ok=True)
        
        log_file = os.path.join(log_dir, f"daily_processor_{datetime.now().strftime('%Y%m%d')}.log")
        
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_file, mode='a', encoding='utf-8'),
                logging.StreamHandler(sys.stdout)
            ]
        )
        self.logger = logging.getLogger(__name__)
    
    def prepare_directories(self):
        """准备目录结构"""
        directories = [
            self.config['output_dir'],
            os.path.join(self.config['output_dir'], 'processed'),
            os.path.join(self.config['output_dir'], 'reports'),
            os.path.join(self.config['output_dir'], 'failed')
        ]
        
        for dir_path in directories:
            os.makedirs(dir_path, exist_ok=True)
    
    def validate_input(self):
        """验证输入目录"""
        if not os.path.exists(self.config['input_dir']):
            self.logger.error(f"输入目录不存在: {self.config['input_dir']}")
            return False
        
        pdf_files = [f for f in os.listdir(self.config['input_dir']) if f.lower().endswith('.pdf')]
        
        if not pdf_files:
            self.logger.warning("输入目录中没有找到PDF文件")
            return False
        
        self.logger.info(f"找到 {len(pdf_files)} 个PDF文件需要处理")
        
        return True
    
    def run_processing_pipeline(self):
        """运行处理管道"""
        # 创建处理管道实例
        pipeline = pdf_processing_pipeline.PDFProcessingPipeline(self.config)
        
        # 传递固定缓存文件路径
        self.config['cache_file'] = self.cache_file
        
        success = pipeline.run_pipeline(
            self.config['input_dir'], 
            os.path.join(self.config['output_dir'], 'processed')
        )
        
        if success:
            # 生成每日报告
            pipeline.generate_summary_report(os.path.join(self.config['output_dir'], 'processed'))
        
        return success, pipeline
    
    def move_processed_files(self):
        """将已处理的PDF文件移动到B文件夹"""
        input_dir = self.config['input_dir']
        # B文件夹路径：input_dir的同级目录，命名为 files_processed
        processed_dir = os.path.join(os.path.dirname(input_dir), "files_processed")
        
        os.makedirs(processed_dir, exist_ok=True)
        
        # 获取已处理成功的文件列表
        processed_md_dir = os.path.join(self.config['output_dir'], 'processed')
        if not os.path.exists(processed_md_dir):
            self.logger.warning(f"处理目录不存在: {processed_md_dir}")
            return
        
        # 获取所有已生成的markdown文件名（不含扩展名）
        processed_files = set()
        for md_file in os.listdir(processed_md_dir):
            if md_file.endswith('.md'):
                # 移除 .md 扩展名
                base_name = md_file[:-3]
                # 移除可能的后缀如 _hybrid
                if '_hybrid' in base_name:
                    base_name = base_name.replace('_hybrid', '')
                elif '_tesseract' in base_name:
                    base_name = base_name.replace('_tesseract', '')
                elif '_liteparse' in base_name:
                    base_name = base_name.replace('_liteparse', '')
                processed_files.add(base_name)
        
        # 移动对应的PDF文件
        moved_count = 0
        for pdf_file in os.listdir(input_dir):
            if pdf_file.lower().endswith('.pdf'):
                pdf_base_name = pdf_file[:-4]  # 移除 .pdf 扩展名
                if pdf_base_name in processed_files:
                    src_path = os.path.join(input_dir, pdf_file)
                    dst_path = os.path.join(processed_dir, pdf_file)
                    
                    try:
                        shutil.move(src_path, dst_path)
                        self.logger.info(f"已移动文件: {pdf_file} -> {processed_dir}")
                        moved_count += 1
                    except Exception as e:
                        self.logger.error(f"移动文件失败 {pdf_file}: {e}")
        
        self.logger.info(f"共移动 {moved_count} 个文件到 {processed_dir}")
    
    def run_ai_analysis(self):
        """运行AI内容分析（只处理新文件，并行处理）"""
        processed_dir = os.path.join(self.config['output_dir'], 'processed')
        summary_list_file = os.path.join(self.config['output_dir'], 'reports', 
                                     f"summary_list_{datetime.now().strftime('%Y%m%d')}.md")
        summary_dir = os.path.join(self.config['output_dir'], 'summaries')
        
        # 设置API Key（从环境变量或.env读取）
        import os as _os
        if not _os.environ.get('OPENAI_API_KEY'):
            try:
                from dotenv import load_dotenv
                load_dotenv()
            except ImportError:
                pass
        
        # 确保summaries目录存在
        _os.makedirs(summary_dir, exist_ok=True)
        
        summarizer = ai_content_summarizer.MarkdownFileSummarizer()
        # 只处理新文件，6个并行进程
        analyses = summarizer.batch_process_markdown_files(
            processed_dir, 
            only_new=True, 
            summaries_dir=summary_dir,
            workers=6
        )
        
        if analyses:
            summarizer.generate_summary_outputs(analyses, summary_list_file, summary_dir)
            self.logger.info(f"总结清单已保存到: {summary_list_file}")
            self.logger.info(f"单文件总结已保存到: {summary_dir}")
            self.logger.info(f"本次共处理 {len(analyses)} 个新文件")
        else:
            self.logger.info("没有新文件需要处理（所有文件已有总结）")
    
    def clean_up(self):
        """清理临时文件"""
        try:
            shutil.rmtree(self.temp_dir)
            self.logger.info("临时目录已清理")
        except Exception as e:
            self.logger.warning(f"清理临时文件失败: {e}")
    
    def run_optimization_analysis(self):
        """运行处理优化分析"""
        processed_dir = os.path.join(self.config['output_dir'], 'processed')
        reports_dir = os.path.join(self.config['output_dir'], 'reports')
        
        if not os.path.exists(processed_dir):
            return
        
        md_files = [f for f in os.listdir(processed_dir) if f.endswith('.md')]
        
        if not md_files:
            return
        
        optimization_report = {
            'total_files': len(md_files),
            'avg_char_count': 0,
            'avg_processing_time': 0,
            'files_per_minute': 0,
            'strategies_used': {},
            'file_sizes': []
        }
        
        char_counts = []
        
        for md_file in md_files:
            file_path = os.path.join(processed_dir, md_file)
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                char_counts.append(len(content))
            except Exception as e:
                self.logger.error(f"读取文件时出错 {md_file}: {e}")
        
        if char_counts:
            optimization_report['avg_char_count'] = sum(char_counts) / len(char_counts)
        
        self.generate_optimization_report(optimization_report, reports_dir)
    
    def generate_optimization_report(self, report, reports_dir):
        """生成优化报告"""
        report_file = os.path.join(reports_dir, f"optimization_report_{datetime.now().strftime('%Y%m%d%H')}.md")
        
        content = f"# 每日处理优化报告 {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"
        
        content += f"## 处理统计\n"
        content += f"- 成功处理文件: {report['total_files']}\n"
        content += f"- 平均字符数: {report['avg_char_count']:.0f}\n"
        
        content += f"\n## 优化建议\n"
        content += f"根据处理结果，建议:\n\n"
        content += f"- 如果平均字符数 < 200，可能需要调整OCR策略\n"
        content += f"- 如果处理时间过长，考虑调整工作进程数\n"
        content += f"- 对于频繁失败的文件，检查源文件质量\n"
        
        with open(report_file, 'w', encoding='utf-8') as f:
            f.write(content)
        
        self.logger.info(f"优化报告已生成: {os.path.basename(report_file)}")
    
    def convert_office_documents(self):
        """将输入目录中的 Office 文档（docx/doc/xlsx/xls/pptx）转换为 Markdown，纳入总结流程

        - .docx/.xlsx/.xls: 使用 markitdown 转换（.xls 依赖 xlrd）
        - .doc:            使用 macOS 自带 textutil 转换（零依赖）
        - .pptx:           markitdown 提取文本 + 图片 Tesseract OCR 合并
        转换成功后源文件移动到 files_processed/
        """
        input_dir = self.config['input_dir']
        processed_dir = os.path.join(self.config['output_dir'], 'processed')
        os.makedirs(processed_dir, exist_ok=True)

        processed_dir_out = os.path.join(os.path.dirname(input_dir), "files_processed")
        os.makedirs(processed_dir_out, exist_ok=True)

        converted_count = 0
        for filename in sorted(os.listdir(input_dir)):
            ext = os.path.splitext(filename)[1].lower()
            if ext not in ('.docx', '.doc', '.xlsx', '.xls', '.pptx'):
                continue
            src_path = os.path.join(input_dir, filename)
            if not os.path.isfile(src_path):
                continue

            try:
                content = self._convert_office_to_markdown(src_path, ext)
                if not content or not content.strip():
                    self.logger.warning(f"转换后内容为空: {filename}，跳过")
                    continue

                md_name = os.path.splitext(filename)[0] + '_office.md'
                md_path = os.path.join(processed_dir, md_name)
                with open(md_path, 'w', encoding='utf-8') as f:
                    f.write(content)

                dst_path = os.path.join(processed_dir_out, filename)
                base, e = os.path.splitext(filename)
                counter = 1
                while os.path.exists(dst_path):
                    dst_path = os.path.join(processed_dir_out, f"{base}_{counter}{e}")
                    counter += 1
                shutil.move(src_path, dst_path)

                self.logger.info(f"已转换并归档: {filename} -> {md_name}")
                converted_count += 1
            except Exception as e:
                self.logger.error(f"转换失败 {filename}: {e}")

        if converted_count > 0:
            self.logger.info(f"共转换 {converted_count} 个 Office 文档")
        else:
            self.logger.info("未检测到 Office 文档")

        return converted_count

    def _convert_office_to_markdown(self, file_path, ext):
        """转换 Office 文档（Word/Excel/PPT）为 Markdown 文本"""
        if ext == '.doc':
            result = subprocess.run(
                ['/usr/bin/textutil', '-convert', 'txt', '-stdout', file_path],
                capture_output=True, text=True, encoding='utf-8', errors='replace'
            )
            if result.returncode != 0:
                raise RuntimeError(f"textutil 转换失败: {result.stderr}")
            return result.stdout
        if ext == '.pptx':
            return self._convert_pptx_to_markdown(file_path)
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content

    def _convert_pptx_to_markdown(self, file_path):
        """转换 PPTX：markitdown 提取文本元素 + python-pptx 提取图片做 Tesseract OCR，合并输出

        文字型 PPT 通过 markitdown 提取标题/文本框/备注；
        图片型 PPT（文字在图片里）通过提取嵌入图片 + OCR 补充文字。
        """
        import hashlib
        from markitdown import MarkItDown
        md = MarkItDown()
        text = md.convert(file_path).text_content

        ocr_parts = []
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import pytesseract
            import PIL.ImageOps
            from PIL import Image

            prs = Presentation(file_path)
            seen = set()
            for slide_idx, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                    try:
                        img = shape.image
                        data = img.blob
                        digest = hashlib.md5(data).hexdigest()
                        if digest in seen:
                            continue
                        seen.add(digest)
                        with tempfile.NamedTemporaryFile(suffix=f".{img.ext}", delete=False) as tmp:
                            tmp.write(data)
                            tmp_path = tmp.name
                        try:
                            pil_img = Image.open(tmp_path)
                            w, h = pil_img.size
                            if w < 80 or h < 80:
                                continue
                            enhanced = PIL.ImageOps.autocontrast(pil_img.convert('L'), cutoff=2)
                            enhanced = enhanced.point(lambda x: 0 if x < 150 else 255, '1')
                            ocr_text = pytesseract.image_to_string(enhanced, lang='chi_sim+eng').strip()
                            if ocr_text:
                                ocr_parts.append(f"=== PPT 第 {slide_idx} 页图片 ===\n{ocr_text}")
                        finally:
                            os.remove(tmp_path)
                    except Exception as e:
                        self.logger.warning(f"PPT 图片 OCR 失败: {e}")
        except Exception as e:
            self.logger.warning(f"PPT 图片提取失败: {e}")

        if ocr_parts:
            text = (text + "\n\n" + "\n\n".join(ocr_parts)).strip()
        return text

    def move_non_pdf_files(self):
        """将输入目录中的非PDF文件移动到 Downloads 目录"""
        move_to_dir = os.path.expanduser('~/Downloads')
        moved_count = 0
        
        for filename in os.listdir(self.config['input_dir']):
            src_path = os.path.join(self.config['input_dir'], filename)
            if not os.path.isfile(src_path):
                continue
            if filename.lower().endswith('.pdf'):
                continue
            
            dst_path = os.path.join(move_to_dir, filename)
            # 处理重名文件
            base, ext = os.path.splitext(filename)
            counter = 1
            while os.path.exists(dst_path):
                dst_path = os.path.join(move_to_dir, f"{base}_{counter}{ext}")
                counter += 1
            
            shutil.move(src_path, dst_path)
            moved_count += 1
            self.logger.info(f"移除非PDF文件: {filename} -> Downloads")
        
        if moved_count > 0:
            self.logger.info(f"共移动 {moved_count} 个非PDF文件到 {move_to_dir}")
        else:
            self.logger.info("未检测到非PDF文件")
    
    def run_all(self):
        """运行完整的每日处理流程"""
        self.logger.info("=== 开始每日PDF处理系统 ===")
        
        try:
            # 准备目录
            self.prepare_directories()
            
            # 转换 Word 文档（docx/doc -> Markdown，纳入总结流程）
            self.convert_office_documents()
            
            # 移除非PDF文件（默认流程）
            self.move_non_pdf_files()
            
            # 验证输入
            if not self.validate_input():
                self.logger.warning("没有找到有效的PDF文件，处理流程将终止")
                return False
            
            # 运行处理管道
            processing_start_time = time.time()
            self.logger.info("正在运行处理管道...")
            pipeline_success, pipeline = self.run_processing_pipeline()
            
            if pipeline_success:
                self.logger.info("处理管道运行成功")
            else:
                self.logger.error("处理管道运行失败")
                return False
            
            # 运行AI内容分析
            if not self.config.get('no_ai', False):
                self.logger.info("正在运行AI内容分析...")
                self.run_ai_analysis()
            else:
                self.logger.info("已跳过AI内容分析")
            
            # 提取多个专题报告
            if not self.config.get('no_ai', False):
                from extract_topic_summary import extract_topic_by_keywords, TOPIC_CONFIGS
                summaries_dir = os.path.join(self.config['output_dir'], 'summaries')
                topics = self.config.get('topics', ['AI'])
                topic_date_str = os.path.basename(os.path.normpath(self.config['output_dir']))
                
                for topic in topics:
                    if topic in TOPIC_CONFIGS:
                        self.logger.info(f"正在提取【{topic}】专题报告...")
                        extract_topic_by_keywords(summaries_dir, TOPIC_CONFIGS[topic], date_str=topic_date_str)
                    else:
                        self.logger.warning(f"不支持的专题: {topic}，跳过")
            
            # 生成每日重点汇总
            if not self.config.get('no_ai', False):
                summaries_dir = os.path.join(self.config['output_dir'], 'summaries')
                reports_dir = os.path.join(self.config['output_dir'], 'reports')
                topic_date_str = os.path.basename(os.path.normpath(self.config['output_dir']))
                highlight_file = os.path.join(reports_dir, f'每日重点汇总_{topic_date_str}.md')
                
                self.logger.info("正在生成每日重点汇总...")
                try:
                    from pdf_ocr_tool.summarizers.financial_summarizer import FinancialResearchSummarizer
                    summarizer = FinancialResearchSummarizer(use_llm=True)
                    summarizer.generate_daily_highlight_report(
                        summaries_dir, highlight_file, date_str=topic_date_str
                    )
                except Exception as e:
                    self.logger.warning(f"生成每日重点汇总失败: {e}")
            
            # 复制每日重点汇总和一句话总结到集合文件夹
            self.logger.info("正在复制汇总文档到集合文件夹...")
            try:
                copied = copy_reports_to_collections(self.config['output_dir'])
                for src, dst in copied:
                    self.logger.info(f"已复制: {os.path.basename(src)} -> {dst}")
                if not copied:
                    self.logger.info("无可复制的汇总文档")
            except Exception as e:
                self.logger.warning(f"复制汇总文档失败: {e}")
            
            # 收集当日汇总文档到 output/汇总/YYYYMMDD/（微信读书导入用）
            self.logger.info("正在收集当日汇总文档到汇总文件夹...")
            try:
                collected = collect_daily_summaries(self.config['output_dir'])
                for src, dst in collected:
                    self.logger.info(f"已收集: {os.path.basename(src)} -> {dst}")
                if not collected:
                    self.logger.info("无可收集的汇总文档")
            except Exception as e:
                self.logger.warning(f"收集汇总文档失败: {e}")
            
            # 运行优化分析
            self.logger.info("正在生成优化报告...")
            self.run_optimization_analysis()
            
            # 移动已处理的文件到B文件夹
            self.logger.info("正在移动已处理的文件...")
            self.move_processed_files()
            
            total_time = time.time() - processing_start_time
            self.logger.info(f"=== 处理完成 ===")
            self.logger.info(f"总处理时间: {total_time:.2f} 秒")
            self.logger.info(f"输出目录: {self.config['output_dir']}")
            
            return True
            
        except Exception as e:
            self.logger.error(f"处理过程中发生错误: {e}")
            import traceback
            self.logger.error(f"错误详情: {traceback.format_exc()}")
            return False
            
        finally:
            self.clean_up()
